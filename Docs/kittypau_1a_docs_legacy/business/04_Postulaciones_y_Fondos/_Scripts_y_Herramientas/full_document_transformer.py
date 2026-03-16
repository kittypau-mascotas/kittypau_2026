import os
import shutil
# Import necessary libraries (these need to be installed via pip)
try:
    from pypdf import PdfReader # For PDF
    from docx import Document # For DOCX
    import openpyxl # For XLSX
    import pandas as pd # For CSV and potentially XLSX
    from pptx import Presentation # For PPTX
except ImportError:
    print("Una o más librerías requeridas no están instaladas. Por favor, ejecuta install_dependencies.ps1")
    exit()

def convert_pdf_to_text(pdf_path, output_path):
    """Converts a PDF file to plain text."""
    try:
        reader = PdfReader(pdf_path)
        text = ""
        for page in reader.pages:
            text += page.extract_text() + "\n"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        return True
    except Exception as e:
        print(f"Error al convertir PDF {pdf_path}: {e}")
        return False

def convert_docx_to_text(docx_path, output_path):
    """Converts a DOCX file to plain text."""
    try:
        document = Document(docx_path)
        text = ""
        for para in document.paragraphs:
            text += para.text + "\n"
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(text)
        return True
    except Exception as e:
        print(f"Error al convertir DOCX {docx_path}: {e}")
        return False

def convert_xlsx_to_csv(xlsx_path, output_path):
    """Converts an XLSX file to CSV (first sheet only)."""
    try:
        df = pd.read_excel(xlsx_path, engine='openpyxl')
        df.to_csv(output_path, index=False, encoding='utf-8')
        return True
    except Exception as e:
        print(f"Error al convertir XLSX {xlsx_path}: {e}")
        return False

def convert_csv_to_md(csv_path, output_path):
    """Converts a CSV file to Markdown table format."""
    try:
        with open(csv_path, 'r', encoding='utf-8') as f:
            content = f.read()
        
        lines = content.strip().split('\n')
        if not lines:
            return False
        
        header = lines[0].split(',')
        md_content = "| " + " | ".join(header) + " |\n"
        md_content += "| " + " | ".join(["---"] * len(header)) + " |\n"
        
        for line in lines[1:]:
            row = line.split(',')
            md_content += "| " + " | ".join(row) + " |\n"
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(md_content)
        return True
    except Exception as e:
        print(f"Error al convertir CSV {csv_path}: {e}")
        return False

def convert_pptx_to_text(pptx_path, output_path):
    """Extracts text from a PPTX file."""
    try:
        prs = Presentation(pptx_path)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    text_runs.append(shape.text)
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write("\n\n".join(text_runs))
        return True
    except Exception as e:
        print(f"Error al convertir PPTX {pptx_path}: {e}")
        return False

def process_incompatible_files(root_source_dir, root_destination_dir):
    incompatible_folder = os.path.join(root_source_dir, "Formato_no_Compatible")
    
    if not os.path.exists(incompatible_folder):
        print(f"Carpeta no encontrada: {incompatible_folder}")
        return

    for dirpath, dirnames, filenames in os.walk(incompatible_folder):
        for filename in filenames:
            source_file_path = os.path.join(dirpath, filename)
            
            # Determinar la ruta relativa desde la carpeta incompatible_folder
            relative_path = os.path.relpath(source_file_path, incompatible_folder)
            
            # Construir el directorio de destino en root_destination_dir
            # Esto asegura que la estructura de carpetas original se mantenga
            target_dir = os.path.join(root_destination_dir, os.path.dirname(relative_path))
            os.makedirs(target_dir, exist_ok=True)

            file_name_without_ext, file_extension = os.path.splitext(filename)
            file_extension = file_extension.lower()
            
            converted = False
            output_file_name = ""

            if file_extension == ".pdf":
                output_file_name = file_name_without_ext + ".txt"
                converted = convert_pdf_to_text(source_file_path, os.path.join(target_dir, output_file_name))
            elif file_extension == ".docx":
                output_file_name = file_name_without_ext + ".txt"
                converted = convert_docx_to_text(source_file_path, os.path.join(target_dir, output_file_name))
            elif file_extension == ".xlsx":
                # Convertir XLSX a CSV primero, luego CSV a MD
                temp_csv_path = os.path.join(target_dir, file_name_without_ext + ".csv")
                if convert_xlsx_to_csv(source_file_path, temp_csv_path):
                    output_file_name = file_name_without_ext + ".md"
                    converted = convert_csv_to_md(temp_csv_path, os.path.join(target_dir, output_file_name))
                    os.remove(temp_csv_path) # Limpiar el CSV temporal
            elif file_extension == ".pptx":
                output_file_name = file_name_without_ext + ".txt"
                converted = convert_pptx_to_text(source_file_path, os.path.join(target_dir, output_file_name))
            
            if converted:
                print(f"Convertido exitosamente {filename} a {output_file_name} en {target_dir}")
            else:
                print(f"No se pudo convertir {filename}. Permanece en {source_file_path}")

if __name__ == "__main__":
    # Define el directorio raíz donde se encuentra la carpeta 'Formato_no_Compatible'
    # y donde se deben colocar los archivos convertidos (la carpeta principal de docs)
    root_docs_path = "D:\\Escritorio\\Proyectos\\KittyPaw\\Kittypaw_1a\\docs\\business\\04_Postulaciones_y_Fondos"
    
    print("Iniciando el proceso de transformación de documentos...")
    process_incompatible_files(root_docs_path, root_docs_path)
    print("Proceso de transformación de documentos finalizado.")
